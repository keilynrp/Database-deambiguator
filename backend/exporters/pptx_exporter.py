"""
PowerPoint exporter — Phase 10 (Artifact Studio)
Generates a branded 16:9 PPTX from report data sections.
Requires python-pptx >= 1.0.2
"""
from __future__ import annotations

import io
from datetime import datetime, timezone
from typing import List, Optional

from sqlalchemy import func
from sqlalchemy.orm import Session

from backend import models
from backend.tenant_access import scope_query_to_org

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#rrggbb' to RGBColor."""
    h = hex_color.lstrip("#")
    if len(h) != 6:
        h = "6366f1"
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def _add_slide(prs: "Presentation") -> any:
    blank_layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(blank_layout)


def _add_header_bar(slide, accent: "RGBColor", width_emu: int, height_emu: int = 914400 // 8):
    """Add a colored top bar to the slide."""
    bar = slide.shapes.add_shape(1, 0, 0, width_emu, height_emu)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size: int = 18, bold: bool = False,
                  color: Optional["RGBColor"] = None, wrap: bool = True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def generate_pptx(
    db: Session,
    domain_id: str,
    sections: List[str],
    title: Optional[str],
    branding: dict,
    org_id: int | None = None,
    manual_sections: list[dict[str, str]] | None = None,
    language: str | None = None,
) -> bytes:
    """
    Build a branded 16:9 PPTX.
    Returns raw bytes of the .pptx file.
    Raises ImportError if python-pptx is not installed.
    """
    if not _PPTX_AVAILABLE:
        raise ImportError("python-pptx is required for PowerPoint export.")

    # Resolve deprecated aliases (e.g. top_brands) to public ids so every slide
    # gate below matches the vocabulary GET /reports/sections actually returns.
    from backend.report_builder import canonical_sections
    sections = canonical_sections(sections)

    accent = _hex_to_rgb(branding.get("accent_color", "#6366f1"))
    platform = branding.get("platform_name", "UKIP")
    footer_text = branding.get("footer_text", "Universal Knowledge Intelligence Platform")
    report_title = title or f"{platform} Report"

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    entities_query = scope_query_to_org(db.query(models.RawEntity), models.RawEntity, org_id)
    if domain_id:
        entities_query = entities_query.filter(models.RawEntity.domain == domain_id)

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    slide = _add_slide(prs)
    # Full-width accent rectangle (top third)
    cover_bar = slide.shapes.add_shape(1, 0, 0, W, int(H * 0.45))
    cover_bar.fill.solid()
    cover_bar.fill.fore_color.rgb = accent
    cover_bar.line.fill.background()
    # Platform name (white, top area)
    _add_text_box(slide, platform, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
                  font_size=20, bold=True, color=RGBColor(255, 255, 255))
    # Report title (white, larger)
    _add_text_box(slide, report_title, Inches(1), Inches(2.0), Inches(11), Inches(1.1),
                  font_size=32, bold=True, color=RGBColor(255, 255, 255))
    # Date
    date_str = datetime.now(timezone.utc).strftime("%B %d, %Y")
    _add_text_box(slide, date_str, Inches(1), Inches(3.0), Inches(6), Inches(0.4),
                  font_size=13, color=RGBColor(220, 220, 255))
    # Domain
    _add_text_box(slide, f"Domain: {domain_id}", Inches(1), Inches(5.5), Inches(8), Inches(0.4),
                  font_size=12, color=RGBColor(100, 100, 120))
    # Footer
    _add_text_box(slide, footer_text, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
                  font_size=10, color=RGBColor(150, 150, 150))

    for manual in manual_sections or []:
        manual_title = (manual.get("title") or "Analyst Note").strip()[:120]
        manual_content = (manual.get("content") or "").strip()
        if not manual_content:
            continue
        slide = _add_slide(prs)
        _add_header_bar(slide, accent, W)
        _add_text_box(slide, manual_title or "Analyst Note", Inches(0.5), Inches(0.05), Inches(10), Inches(0.45),
                      font_size=16, bold=True, color=RGBColor(255, 255, 255))
        _add_text_box(slide, manual_content[:1800], Inches(0.75), Inches(0.9), Inches(11.8), Inches(5.6),
                      font_size=15, color=RGBColor(45, 55, 72))

    # entity_stats, enrichment_coverage and top_secondary_labels used to be
    # hand-built slides here, each issuing its own queries. That made them the
    # last three sections bypassing the shared payload in this format — the same
    # violation 3.3 found in topic_clusters — so they carried no takeaway and no
    # disclosure while the parity map claimed PPTX rendered them. They are in the
    # migrated map below now. Migrating cost no detail: the payload is richer
    # than all three were (four KPI cards rather than two, a Source column, 15
    # rows rather than 10).

    # topic_clusters used to be a hand-built two-column slide here, titled "Top
    # Concepts" and fetching its own top_n=20 while Excel used 50 and HTML 15.
    # It is now in the migrated map below, so all four formats render the same
    # rows from one payload.

    # ── Migrated sections: rendered from the shared section payload ───────────
    # These render via the format-neutral collectors + the shared PPTX renderer,
    # so the section is authored once and appears here without a bespoke slide
    # builder. (unify-report-format-coverage phase 3.) The hand-written slides
    # above are left in place for now; de-duping them onto their collectors is
    # deferred to the cleanup phase. `sections` is already canonicalized above.
    from dataclasses import replace as _replace

    from backend import report_builder
    from backend.i18n.catalog import translate
    from backend.i18n.locale import resolve_report_language
    from backend.reporting.pptx_renderer import render_pptx

    language = resolve_report_language(language)
    migrated_collectors = {
        "entity_stats": report_builder.collect_entity_stats,
        "enrichment_coverage": report_builder.collect_enrichment_coverage,
        "top_secondary_labels": report_builder.collect_top_secondary_labels,
        "impact_projection": report_builder.collect_impact_projection,
        "institutional_benchmark": report_builder.collect_institutional_benchmark,
        "hidden_patterns": report_builder.collect_hidden_patterns,
        "decision_recommendations": report_builder.collect_decision_recommendations,
        "harmonization_log": report_builder.collect_harmonization_log,
        "authority_control": report_builder.collect_authority_control,
        "collaboration_graph": report_builder.collect_collaboration_graph,
        "journal_portfolio": report_builder.collect_journal_portfolio,
        "topic_clusters": report_builder.collect_topic_clusters,
    }
    for section_id, collect in migrated_collectors.items():
        if section_id in sections:
            # A slide heading has no length limit, so it takes the full title
            # rather than Excel's abbreviated sheet name.
            payload = _replace(
                collect(db, domain_id, org_id),
                title=translate(f"report.section.{section_id}", language),
            )
            render_pptx(payload, prs, accent)

    # ── Final slide: Closing ──────────────────────────────────────────────────
    slide = _add_slide(prs)
    closing_bar = slide.shapes.add_shape(1, 0, 0, W, H)
    closing_bar.fill.solid()
    closing_bar.fill.fore_color.rgb = accent
    closing_bar.line.fill.background()
    _add_text_box(slide, platform, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
                  font_size=28, bold=True, color=RGBColor(255, 255, 255))
    _add_text_box(slide, footer_text, Inches(1), Inches(3.5), Inches(11), Inches(0.5),
                  font_size=14, color=RGBColor(220, 220, 255))
    _add_text_box(slide, f"Generated {datetime.now(timezone.utc).strftime('%Y-%m-%d')}",
                  Inches(1), Inches(6.8), Inches(6), Inches(0.4),
                  font_size=10, color=RGBColor(200, 200, 230))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
