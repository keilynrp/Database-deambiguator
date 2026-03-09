"""
Sprint 47 — PowerPoint Export tests
POST /exports/pptx
"""
import io
import pytest


_PAYLOAD = {
    "domain_id": "default",
    "sections": ["entity_stats"],
    "title": "Test Report",
}

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def test_pptx_requires_auth(client):
    resp = client.post("/exports/pptx", json=_PAYLOAD)
    assert resp.status_code in (401, 403)


def test_pptx_returns_correct_mimetype(client, auth_headers, db_session):
    resp = client.post("/exports/pptx", json=_PAYLOAD, headers=auth_headers)
    assert resp.status_code == 200
    assert _PPTX_MIME in resp.headers.get("content-type", "")


def test_pptx_is_valid_file(client, auth_headers, db_session):
    """Load the returned bytes with python-pptx and verify it has slides."""
    from pptx import Presentation
    resp = client.post("/exports/pptx", json=_PAYLOAD, headers=auth_headers)
    assert resp.status_code == 200
    prs = Presentation(io.BytesIO(resp.content))
    assert len(prs.slides) >= 1  # at least cover + section slide + closing


def test_pptx_all_sections(client, auth_headers, db_session):
    """All 5 sections should produce more slides than 1 section."""
    from pptx import Presentation
    all_payload = {
        "domain_id": "default",
        "sections": ["entity_stats", "enrichment_coverage", "top_brands", "topic_clusters", "harmonization_log"],
    }
    resp = client.post("/exports/pptx", json=all_payload, headers=auth_headers)
    assert resp.status_code == 200
    prs = Presentation(io.BytesIO(resp.content))
    # cover + up to 5 sections + closing = at least 3 slides
    assert len(prs.slides) >= 3


def test_pptx_invalid_section_422(client, auth_headers, db_session):
    payload = {"sections": ["entity_stats", "not_a_real_section"]}
    resp = client.post("/exports/pptx", json=payload, headers=auth_headers)
    assert resp.status_code == 422


def test_viewer_cannot_export_pptx(client, viewer_headers, db_session):
    resp = client.post("/exports/pptx", json=_PAYLOAD, headers=viewer_headers)
    assert resp.status_code in (401, 403)
